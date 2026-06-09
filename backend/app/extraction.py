"""Document extraction: turn an uploaded file into raw text + structured fields.

All parsers are open-source and free. Image OCR uses Tesseract (via pytesseract)
when the binary is available; otherwise the file is still stored and the rest of
the pipeline keeps working (raw_text will be empty for that file).
"""
from __future__ import annotations

import io
import os
import re
from collections import Counter

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}


def detect_kind(filename: str, content_type: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf" or "pdf" in content_type:
        return "pdf"
    if ext in IMAGE_EXTS or content_type.startswith("image/"):
        return "image"
    if ext in {".docx", ".doc"}:
        return "docx"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    if ext in {".xlsx", ".xls", ".csv"}:
        return "sheet"
    if ext in {".txt", ".md", ".rtf"}:
        return "text"
    return "other"


# ---------------------------------------------------------------------------
# Per-format text extractors
# ---------------------------------------------------------------------------
def _ocr_image_bytes(data: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return ""
    try:
        img = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(img, lang="spa+eng")
    except pytesseract.TesseractNotFoundError:
        return ""
    except Exception:
        # Fall back to default language data if spa+eng is unavailable.
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(io.BytesIO(data))
            return pytesseract.image_to_string(img)
        except Exception:
            return ""


def _extract_pdf(data: bytes) -> str:
    text_parts: list[str] = []
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                text_parts.append(page.extract_text() or "")
    except Exception:
        pass
    text = "\n".join(text_parts).strip()
    if text:
        return text
    # Scanned PDF: render pages and OCR them.
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=data, filetype="pdf")
        ocr_parts: list[str] = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            ocr_parts.append(_ocr_image_bytes(pix.tobytes("png")))
        return "\n".join(ocr_parts).strip()
    except Exception:
        return ""


def _extract_docx(data: bytes) -> str:
    try:
        import docx

        document = docx.Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(
                            "\t".join(cell.text for cell in row.cells)
                        )
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_sheet(filename: str, data: bytes) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".csv":
        try:
            return data.decode("utf-8", errors="replace").strip()
        except Exception:
            return ""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts).strip()
    except Exception:
        return ""


def extract_text(filename: str, content_type: str, data: bytes) -> tuple[str, str]:
    """Return (kind, raw_text)."""
    kind = detect_kind(filename, content_type)
    if kind == "pdf":
        return kind, _extract_pdf(data)
    if kind == "image":
        return kind, _ocr_image_bytes(data)
    if kind == "docx":
        return kind, _extract_docx(data)
    if kind == "pptx":
        return kind, _extract_pptx(data)
    if kind == "sheet":
        return kind, _extract_sheet(filename, data)
    if kind == "text":
        try:
            return kind, data.decode("utf-8", errors="replace").strip()
        except Exception:
            return kind, ""
    return kind, ""


# ---------------------------------------------------------------------------
# Structured field detection (the "smart" part)
# ---------------------------------------------------------------------------
EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}")
PHONE_RE = re.compile(r"(?:\+?\d{1,3}[\s\-.]?)?(?:\(?\d{2,4}\)?[\s\-.]?)\d{3,4}[\s\-.]?\d{3,4}")
DATE_RE = re.compile(
    r"\b(\d{1,2}[/\-.]\d{1,2}[/\-.]\d{2,4}|\d{4}[/\-.]\d{1,2}[/\-.]\d{1,2})\b"
)
RIF_RE = re.compile(r"\b[JGVEP]-?\d{8,9}-?\d?\b", re.IGNORECASE)
# Amounts like 1.234,56 / 1,234.56 / 1234.56 optionally prefixed by a currency symbol.
AMOUNT_RE = re.compile(
    r"(Bs\.?|VES|USD|US\$|\$|€)?\s?(\d{1,3}(?:[.,]\d{3})*(?:[.,]\d{2})|\d+[.,]\d{2})",
    re.IGNORECASE,
)

CURRENCY_MAP = {
    "$": "USD",
    "us$": "USD",
    "usd": "USD",
    "€": "EUR",
    "bs": "VES",
    "bs.": "VES",
    "ves": "VES",
}

SPANISH_STOPWORDS = {
    "de", "la", "el", "en", "y", "a", "los", "las", "del", "que", "se", "por",
    "con", "para", "una", "uno", "un", "es", "su", "al", "lo", "como", "mas",
    "o", "the", "of", "to", "and", "in", "for", "on", "no", "si",
}


def _parse_amount(num: str) -> float | None:
    s = num.strip()
    if "," in s and "." in s:
        # Decide decimal separator by last occurrence.
        if s.rfind(",") > s.rfind("."):
            s = s.replace(".", "").replace(",", ".")
        else:
            s = s.replace(",", "")
    elif "," in s:
        # Treat comma as decimal if it looks like cents, else thousands.
        if re.search(r",\d{2}$", s):
            s = s.replace(".", "").replace(",", ".")
        else:
            s = s.replace(",", "")
    try:
        return float(s)
    except ValueError:
        return None


def detect_fields(text: str) -> dict:
    if not text:
        return {
            "emails": [], "phones": [], "dates": [], "rifs": [],
            "amounts": [], "total_amount": 0.0, "currencies": {},
            "keywords": [],
        }

    emails = sorted(set(EMAIL_RE.findall(text)))
    rifs = sorted({m.upper() for m in RIF_RE.findall(text)})
    dates = sorted(set(DATE_RE.findall(text)))[:50]

    phones = []
    for m in PHONE_RE.findall(text):
        digits = re.sub(r"\D", "", m)
        if 7 <= len(digits) <= 15:
            phones.append(m.strip())
    phones = sorted(set(phones))[:50]

    amounts: list[dict] = []
    currencies: Counter[str] = Counter()
    total = 0.0
    for sym, num in AMOUNT_RE.findall(text):
        value = _parse_amount(num)
        if value is None:
            continue
        cur = CURRENCY_MAP.get(sym.strip().lower(), "") if sym else ""
        amounts.append({"value": value, "currency": cur, "raw": (sym + num).strip()})
        total += value
        if cur:
            currencies[cur] += value
    amounts = amounts[:200]

    # Keyword frequency (simple, language-light).
    words = re.findall(r"[A-Za-zÁÉÍÓÚáéíóúÑñ]{4,}", text.lower())
    counter = Counter(w for w in words if w not in SPANISH_STOPWORDS)
    keywords = counter.most_common(15)

    return {
        "emails": emails[:50],
        "phones": phones,
        "dates": dates,
        "rifs": rifs[:50],
        "amounts": amounts,
        "total_amount": round(total, 2),
        "currencies": {k: round(v, 2) for k, v in currencies.items()},
        "keywords": keywords,
    }
